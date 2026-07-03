"""Generate EdCounselor project presentation (PPTX)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PRIMARY = RGBColor(79, 70, 229)
DARK = RGBColor(15, 23, 42)
GRAY = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(238, 242, 255)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    box2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8.4), Inches(1))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(224, 231, 255)
    p2.alignment = PP_ALIGN.CENTER


def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.4), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, bullets, subtext=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)

    if subtext:
        p = tf.add_paragraph()
        p.text = subtext
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.space_before = Pt(16)


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    for col, (ctitle, items, left) in enumerate(
        [(left_title, left_items, 0.5), (right_title, right_items, 5.2)]
    ):
        box = slide.shapes.add_textbox(Inches(left), Inches(1.4), Inches(4.3), Inches(5))
        tf = box.text_frame
        tf.word_wrap = True
        h = tf.paragraphs[0]
        h.text = ctitle
        h.font.size = Pt(20)
        h.font.bold = True
        h.font.color.rgb = PRIMARY
        h.space_after = Pt(12)
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = DARK
            p.space_after = Pt(6)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "EdCounselor",
        "3-Tier Student Portal\nReact.js • Python Serverless • PostgreSQL",
    )

    add_content_slide(
        prs,
        "Project Overview",
        [
            "Purpose: Student enrollment and digital ID card portal",
            "Users: Students (register, enroll, view ID) + Admin (dashboard)",
            "Live URL: task-edcounselor.vercel.app",
            "Stack: React frontend, Python Flask API, Supabase PostgreSQL",
            "Pattern: REST API calls between frontend and backend",
        ],
    )

    add_section_slide(prs, "Architecture")

    add_content_slide(
        prs,
        "3-Tier Architecture",
        [
            "Tier 1 — Presentation: React.js (UI in browser)",
            "Tier 2 — Application: Python Flask serverless API on Vercel",
            "Tier 3 — Data: PostgreSQL database on Supabase",
            "",
            "Browser → HTTP/JSON → Python API → SQL → Database",
            "Database is never exposed directly to the frontend",
        ],
    )

    add_two_column_slide(
        prs,
        "Platforms: Vercel & Supabase",
        "Vercel (Hosting)",
        [
            "Hosts React static build",
            "Runs Python as serverless functions",
            "Auto-deploy from GitHub push",
            "Free HTTPS + live URL",
            "Environment variables for secrets",
        ],
        "Supabase (Database)",
        [
            "Managed PostgreSQL in the cloud",
            "Web UI for tables & data",
            "Connection pooler for serverless",
            "SSL-encrypted connections",
            "Free tier for student projects",
        ],
    )

    add_section_slide(prs, "Database Layer")

    add_content_slide(
        prs,
        "Database Schema (PostgreSQL)",
        [
            "Table 1: userdata",
            "  • email (PK), password_hash, first_name, last_name, created_at",
            "",
            "Table 2: student_profiles",
            "  • profile_id (PK), student_email (FK → userdata.email)",
            "  • institution, DOB, gender, roll_number, class, section, parent_contact",
            "",
            "Relationship: One user → zero or one profile (Foreign Key)",
        ],
    )

    add_content_slide(
        prs,
        "Key SQL Operations",
        [
            "Register: INSERT INTO userdata (bcrypt hashed password)",
            "Login: SELECT password_hash + bcrypt verification",
            "Create Profile: INSERT INTO student_profiles (FK check)",
            "View Profile: JOIN userdata + student_profiles",
            "Admin List: LEFT JOIN (shows users without profiles too)",
            "Security: Parameterized queries (%s) prevent SQL injection",
        ],
    )

    add_section_slide(prs, "Backend API")

    add_content_slide(
        prs,
        "Python Backend — api/index.py",
        [
            "Framework: Flask 3.0 + psycopg2 + bcrypt",
            "Deployed as Vercel serverless function (handler = app)",
            "",
            "GET  /api/health          — Database connection check",
            "POST /api/register        — User sign up",
            "POST /api/login           — User authentication",
            "POST /api/create-profile  — Save student enrollment",
            "GET  /api/get-profile     — Fetch ID card data",
            "GET  /api/get-all-profiles — Admin dashboard data",
        ],
    )

    add_content_slide(
        prs,
        "API: Register & Login",
        [
            "Register (POST /api/register):",
            "  Input: firstName, lastName, email, password",
            "  Process: Validate → bcrypt hash → INSERT userdata",
            "  Output: { email, hasProfile: false, role: student }",
            "",
            "Login (POST /api/login):",
            "  Input: email, password",
            "  Process: SELECT user → verify hash → check profile exists",
            "  Output: { email, role, hasProfile }",
        ],
    )

    add_section_slide(prs, "Frontend")

    add_content_slide(
        prs,
        "React Frontend Modules",
        [
            "main.jsx — App entry point, mounts React",
            "App.jsx — View controller (auth / create / view / admin)",
            "Auth.jsx — Login & Sign Up page",
            "ProfileCreate.jsx — Student enrollment form",
            "ProfileView.jsx — Digital student ID card",
            "AdminDashboard.jsx — Admin table + statistics",
            "index.css — Unified styling (cards, forms, tables)",
        ],
    )

    add_content_slide(
        prs,
        "Frontend ↔ Backend Connection",
        [
            "React uses fetch() with JSON over HTTP",
            "Same domain: /api/* routed to Python by vercel.json",
            "Example: fetch('/api/login', { method: 'POST', body: JSON })",
            "",
            "No direct database access from browser",
            "DATABASE_URL stored only in Vercel env (not in frontend code)",
            "Responses drive navigation: hasProfile → create or view page",
        ],
    )

    add_section_slide(prs, "User Flows")

    add_content_slide(
        prs,
        "Flow: New Student",
        [
            "1. Open site → Sign Up (Auth.jsx)",
            "2. POST /api/register → row in userdata",
            "3. Redirect to ProfileCreate (no profile yet)",
            "4. Fill enrollment form → POST /api/create-profile",
            "5. Row in student_profiles (FK to userdata)",
            "6. Redirect to ProfileView → GET /api/get-profile",
            "7. Digital ID card displayed",
        ],
    )

    add_content_slide(
        prs,
        "Flow: Returning Student & Admin",
        [
            "Returning student:",
            "  Login → hasProfile: true → ProfileView directly",
            "  Login → hasProfile: false → ProfileCreate first",
            "",
            "Admin:",
            "  Login with admin email → role: admin",
            "  AdminDashboard → GET /api/get-all-profiles",
            "  Table: name, email, institution, roll, class, parent contact",
        ],
    )

    add_section_slide(prs, "Security & Deployment")

    add_two_column_slide(
        prs,
        "Security & Deployment",
        "Security",
        [
            "bcrypt password hashing",
            "Parameterized SQL queries",
            "SSL to database (sslmode=require)",
            "Secrets in Vercel env vars",
            "HTTPS on Vercel automatically",
        ],
        "Deployment",
        [
            "GitHub → Vercel auto-deploy",
            "vercel.json routes /api/* to Python",
            "frontend/ built with Vite",
            "requirements.txt for Python deps",
            "Health check: /api/health",
        ],
    )

    add_content_slide(
        prs,
        "Assessment Requirements Met",
        [
            "✅ React.js frontend",
            "✅ Python serverless backend",
            "✅ PostgreSQL database (Supabase)",
            "✅ REST API calls (fetch + JSON)",
            "✅ Login page",
            "✅ Profile creation page",
            "✅ Profile view page",
            "➕ Bonus: Admin dashboard, styled UI, health endpoint",
        ],
    )

    add_title_slide(
        prs,
        "Thank You",
        "EdCounselor — Student Portal Demo\nQuestions?",
    )

    out = "EdCounselor_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
